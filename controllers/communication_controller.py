# controllers/communication_controller.py
# Importa las funciones de conexión directamente de db.connection
from db.connection import create_connection, close_connection
from pptx import Presentation
from pptx.util import Inches
import os
import re # Necesario para el _on_project_select en la vista si parsea la cadena

class CommunicationController:
    def __init__(self):
        # Ya no necesitamos un objeto db_manager aquí,
        # llamaremos a create_connection() y close_connection() directamente.
        pass

    def get_periods(self):
        """Obtiene una lista de todos los períodos."""
        conn = create_connection()
        if conn is None:
            return None, "Error: No se pudo establecer conexión con la base de datos."
        cursor = conn.cursor(dictionary=True) # Usar dictionary=True para obtener resultados como diccionarios
        try:
            cursor.execute("SELECT id_periodo, nombre_periodo, fecha_inicio, fecha_fin FROM periodos ORDER BY fecha_inicio DESC")
            periods = cursor.fetchall()
            return periods, None
        except Exception as e:
            return None, f"Error al obtener períodos: {e}"
        finally:
            if cursor:
                cursor.close()
            close_connection(conn)

    def get_all_eligible_recipients(self):
        """
        Obtiene todos los usuarios y participantes con email válido.
        Unifica la lista de usuarios y participantes.
        """
        all_recipients = []
        conn = create_connection()
        if conn is None:
            return None, "Error: No se pudo establecer conexión con la base de datos."
        cursor = conn.cursor(dictionary=True)
        try:
            # Obtener usuarios (con rol de estudiante o profesor si es aplicable, o todos los que tienen correo_electronico)
            # ¡CORREGIDO! Cambiado 'email' por 'correo_electronico' en la SELECT
            cursor.execute("SELECT id_usuario, nombre_completo, correo_electronico, rol FROM usuarios WHERE correo_electronico IS NOT NULL AND correo_electronico != ''")
            users = [{"id": row['id_usuario'], "nombre_completo": row['nombre_completo'], "email": row['correo_electronico'], "tipo": row['rol']} for row in cursor.fetchall()]
            all_recipients.extend(users)

            # Obtener participantes
            # ¡CORREGIDO! Cambiado 'email' por 'correo_electronico' en la SELECT
            cursor.execute("SELECT id_participante, nombre, apellido, correo_electronico FROM participantes WHERE correo_electronico IS NOT NULL AND correo_electronico != ''")
            # ¡CORREGIDO! Cambiado 'email' por 'correo_electronico' para acceder a la clave del diccionario
            participants = [{"id": row['id_participante'], "nombre_completo": f"{row['nombre']} {row['apellido']}", "email": row['correo_electronico'], "tipo": "Participante"} for row in cursor.fetchall()]
            all_recipients.extend(participants)

            # Eliminar duplicados si un mismo email aparece como usuario y participante
            unique_recipients = {r['email']: r for r in all_recipients}.values()
            
            return list(unique_recipients), None
        except Exception as e:
            return None, f"Error al obtener los destinatarios para emails: {e}" # Mensaje de error ajustado para ser más específico
        finally:
            if cursor:
                cursor.close()
            close_connection(conn)

    def get_participants_by_period(self, period_id):
        """
        Obtiene los participantes asociados a proyectos dentro de un período específico.
        """
        conn = create_connection()
        if conn is None:
            return None, "Error: No se pudo establecer conexión con la base de datos."
        cursor = conn.cursor(dictionary=True)
        try:
            query = """
            SELECT DISTINCT
                p.id_participante,
                p.nombre,
                p.apellido,
                p.cedula,
                p.correo_electronico AS email -- Usar el nombre de columna correcto 'correo_electronico'
            FROM
                participantes p
            JOIN
                proyectos_participantes pp ON p.id_participante = pp.id_participante
            JOIN
                proyectos pr ON pp.id_proyecto = pr.id_proyecto
            WHERE
                pr.id_periodo = %s -- Cambiado de ? a %s
            ORDER BY
                p.apellido, p.nombre;
            """
            cursor.execute(query, (period_id,))
            participants = [{
                "id": row['id_participante'],
                "nombre_completo": f"{row['nombre']} {row['apellido']}",
                "cedula": row['cedula'],
                "email": row['email'],
                "tipo": "Participante"
            } for row in cursor.fetchall()]
            return participants, None
        except Exception as e:
            return None, f"Error al obtener participantes por período: {e}"
        finally:
            if cursor:
                cursor.close()
            close_connection(conn)

    def get_projects_for_certificates(self):
        """
        Obtiene los proyectos con la información necesaria para los certificados.
        Incluye nombre del proyecto, descripción, periodo y participantes.
        """
        conn = create_connection()
        if conn is None:
            return None, "Error: No se pudo establecer conexión con la base de datos."
        cursor = conn.cursor(dictionary=True)
        try:
            query = """
            SELECT
                p.id_proyecto,
                p.nombre_proyecto,
                p.descripcion,
                pe.nombre_periodo,
                -- CAMBIO APLICADO AQUÍ: Usar CONCAT() para MySQL/MariaDB
                GROUP_CONCAT(CONCAT(pa.nombre, ' ', pa.apellido, ' (CI: ', pa.cedula, ' - ', COALESCE(pa.correo_electronico, 'N/A'), ')') SEPARATOR '; ') AS participantes_info
            FROM
                proyectos p
            LEFT JOIN
                periodos pe ON p.id_periodo = pe.id_periodo
            LEFT JOIN
                proyectos_participantes pp ON p.id_proyecto = pp.id_proyecto
            LEFT JOIN
                participantes pa ON pp.id_participante = pa.id_participante
            GROUP BY
                p.id_proyecto, p.nombre_proyecto, p.descripcion, pe.nombre_periodo
            ORDER BY
                pe.nombre_periodo DESC, p.nombre_proyecto ASC;
            """
            cursor.execute(query)
            projects_data = cursor.fetchall()
            
            for project in projects_data:
                if project['participantes_info'] is None:
                    project['participantes_info'] = "N/A"
            
            return projects_data, None
        except Exception as e:
            return None, f"Error al obtener proyectos para certificados: {e}"
        finally:
            if cursor:
                cursor.close()
            close_connection(conn)

    def generate_certificate(self, participant_name, participant_ci, project_name, event_date="16 de enero del 2025", template_path="Formato.pptx", output_dir="certificados_generados"):
        """
        Genera un certificado personalizado a partir de la plantilla.

        Args:
            participant_name (str): Nombre completo del participante.
            participant_ci (str): Cédula de identidad del participante.
            project_name (str): Nombre del proyecto en el que participó.
            event_date (str): La fecha del evento. (Por ahora fijo, puede hacerse dinámico)
            template_path (str): Ruta al archivo de la plantilla .pptx.
            output_dir (str): Directorio donde se guardarán los certificados generados.

        Returns:
            tuple: (ruta_del_certificado_generado, error_mensaje)
        """
        try:
            if not os.path.exists(template_path):
                return None, f"La plantilla no se encontró en: {template_path}"

            prs = Presentation(template_path)
            slide = prs.slides[0] # Asumiendo que el certificado está en la primera diapositiva

            # Buscar la forma que contiene el texto del participante y el proyecto
            # Basado en la salida del script, es el 'object 3'
            target_shape = None
            for shape in slide.shapes:
                # Comprobar si tiene text_frame y si el nombre es 'object 3'
                if shape.has_text_frame and shape.name == 'object 3':
                    target_shape = shape
                    break
            
            if not target_shape:
                return None, "No se encontró la forma 'object 3' en la plantilla para insertar el texto."

            text_frame = target_shape.text_frame
            
            # Limpiamos todo el contenido existente del text_frame
            # Se eliminan todos los párrafos excepto el primero, y luego se limpia el primero.
            # Esto es más robusto si la forma tiene múltiples párrafos por defecto.
            while len(text_frame.paragraphs) > 1:
                p = text_frame.paragraphs[-1]._element
                p.getparent().remove(p)
            
            # Limpiar el texto del único párrafo restante (o el primero si estaba vacío)
            p = text_frame.paragraphs[0]
            p.text = "" # Borra el texto existente

            # Aquí construimos el texto dinámicamente con saltos de línea para el formato.
            # Puedes ajustar el formato según necesites nuevas líneas, etc.
            # Importante: Esto asume que el text_frame se expandirá para contener el texto.
            # Si el text_frame tiene un tamaño fijo, el texto podría desbordarse.
            
            # Primer párrafo: "Que se otorga a:"
            p.text = "Que se otorga a:"
            
            # Segundo párrafo: Nombre del participante (posiblemente con un estilo diferente)
            new_p_name = text_frame.add_paragraph()
            new_p_name.text = f"{participant_name}"
            # Opcional: Si quieres un estilo diferente para el nombre, puedes aplicarlo aquí
            # from pptx.util import Pt
            # new_p_name.font.size = Pt(24) # Ejemplo: tamaño de fuente más grande
            # new_p_name.font.bold = True # Ejemplo: negrita

            # Tercer párrafo: Cédula de Identidad
            new_p_ci = text_frame.add_paragraph()
            new_p_ci.text = f"C.I. {participant_ci}"

            # Cuarto párrafo: Descripción del proyecto y fecha
            new_p_project = text_frame.add_paragraph()
            new_p_project.text = f"Por haber participado en la 4ta Expoferia de la Escuela de Ingeniería con el proyecto “{project_name}”, realizado el {event_date}."

            # Crear el directorio de salida si no existe
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)

            # Sanitizar el nombre de archivo para evitar caracteres inválidos
            safe_participant_name = re.sub(r'[\\/*?:"<>|]', '', participant_name).replace(' ', '_')
            safe_project_name = re.sub(r'[\\/*?:"<>|]', '', project_name).replace(' ', '_')

            output_filename = f"Certificado_{safe_participant_name}_{safe_project_name}.pptx"
            output_path = os.path.join(output_dir, output_filename)
            prs.save(output_path)
            
            return output_path, None

        except Exception as e:
            return None, f"Error al generar el certificado: {e}"